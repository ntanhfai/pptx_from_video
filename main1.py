import os
import shutil
from os.path import dirname, join

import cv2
import numpy as np
from tqdm import tqdm
import matplotlib.pyplot as plt


def plotLists(*args):
    for k, mList in enumerate(args):
        # Tạo biểu đồ
        plt.plot(mList, label=f'List {k}', marker='*')

    # Thêm tiêu đề và nhãn
    plt.title('Plot of Lists')
    plt.xlabel('Index')
    plt.ylabel('Value')

    # Thêm chú thích
    plt.legend()

    # Hiển thị biểu đồ
    plt.show()


def is_new_slide(prev_frame, current_frame):
    prev_gray = cv2.cvtColor(prev_frame, cv2.COLOR_BGR2GRAY)
    current_gray = cv2.cvtColor(current_frame, cv2.COLOR_BGR2GRAY)
    diff = cv2.absdiff(prev_gray, current_gray)
    _, diff_threshold = cv2.threshold(diff, 30, 255, cv2.THRESH_BINARY)
    similarity = np.sum(diff_threshold == 255) / (diff_threshold.shape[0] * diff_threshold.shape[1])
    # print("curr similar=", similarity)
    return similarity


def sliding_window_max_to_zero(nums, k):
    result = [0] * len(nums)

    for i in range(len(nums) - k + 1):
        window = nums[i:i + k]
        max_value = max(window)
        max_index = window.index(max_value)
        for j in range(i, i + k):
            if j != i + max_index:
                result[j] = 0
            else:
                result[j] = max_value

    return result


def reverse_list(lst):
    return lst[::-1]


def find_zeros_between_nonzeros(lst, threshold=0.018):
    result = []
    start_index = None

    for i, num in enumerate(lst):
        if num > threshold:
            if start_index is not None:
                # Tìm vị trí số 0 giữa hai số > 0
                zero_index = (start_index + i) // 2
                result.append(zero_index)
                start_index = None
        elif start_index is None:
            # Bắt đầu của khoảng 0
            start_index = i

    return result


# tôi có 1 list các số, tôi muốn đọc 1 video, và chỉ lấy vài frame có số thứ tự trong list thôi:
import cv2


def extract_frames_from_video(video_path, frame_numbers):
    cap = cv2.VideoCapture(video_path)
    frames = []

    if not cap.isOpened():
        print("Error: Could not open video.")
        return frames

    frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))

    for frame_number in tqdm(frame_numbers):
        # Kiểm tra xem frame_number có hợp lệ không
        if 0 <= frame_number < frame_count:
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_number)
            ret, frame = cap.read()
            if ret:
                frames.append(frame)
        else:
            print(f"Error: Frame number {frame_number} is out of range.")

    cap.release()
    return frames


def extract_frames_from_folder(FIS, frame_numbers):
    N = len(FIS)
    frames = []
    for frame_num in tqdm(frame_numbers):
        if 0 <= frame_num < N:
            frames.append(FIS[frame_num])
        else:
            print(f"Error: Frame number {frame_num} is out of range.")
    return frames


# # Đường dẫn đến video
# video_path = "path/to/your/video.mp4"
#
# # List các số thứ tự của frame bạn muốn trích xuất
# frame_numbers = [10, 20, 30, 40, 50]
#
# # Trích xuất các frame từ video
# extracted_frames = extract_frames_from_video(video_path, frame_numbers)
#
# # Hiển thị các frame
# for i, frame in enumerate(extracted_frames):
#     cv2.imshow(f"Frame {frame_numbers[i]}", frame)
#     cv2.waitKey(0)  # Chờ người dùng nhấn phím bất kỳ để tiếp tục
#
# cv2.destroyAllWindows()


# # Example:
# lst = [0, 0, 1, 0, 0, 0, 0, 2, 0, 0, 0, 0, 0, 3, 0, 0, 0]
# print(find_zeros_between_nonzeros(lst))  # Output: [3, 8, 15]


# Example:
# nums = [1, 3, -1, -3, 5, 3, 6, 7]
# k = 3
# print(sliding_window_max_to_zero(nums, k))  # Output: [0, 3, 0, 0, 0, 0, 0, 7]
def extract_slides_from_video(video_file, maxFrames=1000):
    cap = cv2.VideoCapture(video_file)
    if not cap.isOpened():
        print("Error: Cannot open video file.")
        return

    frame_count = 0
    ret, frame = cap.read()
    nFrame = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    print(('nFrame:', nFrame))
    # nFrame = 1000
    difflist = []
    prev_frame = frame
    for cnt in tqdm(range(nFrame)):
        ret, frame = cap.read()
        if not ret:
            break
        similarity = is_new_slide(prev_frame, frame)
        difflist.append(similarity)
        prev_frame = frame

    cap.release()
    cv2.destroyAllWindows()

    # ============= filter
    vl = 0
    difflist1_0 = []
    for k in difflist:
        vl = vl * 0.5 + k * 0.5
        difflist1_0.append(vl)
    vl = 0
    difflist1 = []
    for k in difflist1_0:
        vl = vl * 0.5 + k * 0.5
        difflist1.append(vl)

    # ============ Max

    difflist2 = sliding_window_max_to_zero(difflist1, 10)  # Output: [0, 0, 3, 0, 5, 0, 6, 7]
    difflist3 = reverse_list(difflist2)
    difflist3 = sliding_window_max_to_zero(difflist3, 10)  # Output: [0, 0, 3, 0, 5, 0, 6, 7]
    difflist4 = [1] + reverse_list(difflist3) + [1]

    Picklist = find_zeros_between_nonzeros(difflist4, threshold=0)

    retFrames = extract_frames_from_video(video_file, Picklist)
    os.makedirs('Output', exist_ok=True)
    for frame_count, frame in tqdm(enumerate(retFrames)):
        img_path = f'Output/temp{frame_count:>08}.jpg'
        cv2.imwrite(img_path, frame)
    pass
    # Example usage


def fnFIS(path, ext=('.jpg',)):
    ret = []
    for D, _, F in os.walk(path):
        for fn in F:
            if fn.endswith(ext):
                ret.append(join(D, fn).replace('\\', '/'))
    return ret


def extract_slides_from_folders(image_folders, Output1='Output1'):
    FIS = fnFIS(image_folders)
    FIS.sort()
    nFrame = len(FIS)
    print('nImages:', nFrame)
    # nFrame = 1000
    difflist = [['', 0]]
    prev_frame = cv2.imread(FIS[0])
    prev_path = FIS[0]
    for cnt, imgPath in tqdm(enumerate(FIS)):
        frame = cv2.imread(imgPath)

        similarity = is_new_slide(prev_frame, frame)
        difflist.append([prev_path, similarity])
        prev_frame = frame
        prev_path = imgPath

    retFrames = [frmNo for frmNo, x in difflist if x >= 0.002]

    # retFrames = extract_frames_from_folder(FIS, Picklist)
    os.makedirs(Output1, exist_ok=True)
    for frame_count, frameOrg in tqdm(enumerate(retFrames)):
        img_path = frameOrg.replace(Output, Output1)
        shutil.copy(frameOrg, img_path)
    pass
    # Example usage


RUN = [3]
Output = 'Output'
Output1 = 'Output1'

if 1 in RUN:
    # extract video frames with difference
    video_file = 'Data/Introduction to Generative AI.mp4'
    extract_slides_from_video(video_file)
    print('Done 1')
if 2 in RUN:
    # Re-pickup not same images
    extract_slides_from_folders(image_folders=Output, Output1=Output1)
    print('Done 2')
    pass
if 3 in RUN:
    # Create pptx
    FIS = fnFIS(Output1)
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(r"taTemplate.pptx")
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for img_path in FIS:
        slide = prs.slides.add_slide(slide_layout=prs.slide_layouts[2])  # Use blank layout for each slide, open pptx layout to see the orders.
        slide.shapes.add_picture(img_path, top=Inches(1.08), left=Inches(1.37), width=Inches(11.37), height=Inches(6))
        slide.shapes.title.text = "Introduction to GenAI & LLM"

    output_pptx = "ta-FII-AI-Department.pptx"
    prs.save(output_pptx)
    os.startfile(output_pptx)
    print('Done 3')
