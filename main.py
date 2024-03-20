import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches


def is_new_slide(prev_frame, current_frame, threshold=0.9):
    prev_gray = cv2.cvtColor(prev_frame, cv2.COLOR_BGR2GRAY)
    current_gray = cv2.cvtColor(current_frame, cv2.COLOR_BGR2GRAY)
    diff = cv2.absdiff(prev_gray, current_gray)
    _, diff_threshold = cv2.threshold(diff, 30, 255, cv2.THRESH_BINARY)
    similarity = np.sum(diff_threshold == 255) / (diff_threshold.shape[0] * diff_threshold.shape[1])
    # print("curr similar=", similarity)
    return similarity


def extract_slides_from_video(video_file, output_pptx):
    cap = cv2.VideoCapture(video_file)
    if not cap.isOpened():
        print("Error: Cannot open video file.")
        return

    prs = Presentation()

    frame_count = 0
    ret, frame = cap.read()
    prev_frame = []
    N = 10
    rnd_cnt = 0
    prev_frame = [frame for _ in range(N)]

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        rnd_cnt = (rnd_cnt + N + 1) % N
        check_Cnt = (rnd_cnt + N + 1) % N
        similarity = is_new_slide(prev_frame[check_Cnt], frame)
        print(f"Processed frame {frame_count}, similar: {similarity}")
        if similarity > 0.1:
            # slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use blank layout for each slide

            # Convert frame to bytes
            _, img_encoded = cv2.imencode('.jpg', frame)
            img_bytes = img_encoded.tobytes()

            # Add image to slide
            img_path = f'Output/temp{frame_count:>08}.jpg'
            with open(img_path, 'wb') as f:
                f.write(img_bytes)
            # slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10))

        frame_count += 1

        prev_frame[rnd_cnt] = frame

    cap.release()
    cv2.destroyAllWindows()

    prs.save(output_pptx)
    print(f"Presentation saved as {output_pptx}")


# Example usage
video_file = 'Data/Introduction to Generative AI.mp4'
output_pptx = 'Output/output_presentation.pptx'
extract_slides_from_video(video_file, output_pptx)
